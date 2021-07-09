from pptx import Presentation
import os
import subprocess, sys

#For picture
from pptx.util import Inches

#For Autoshapes

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR


pr1 = Presentation()

slide1_register = pr1.slide_layouts[0]


print(pr1.slide_layouts)
"Add initial slide to presentation"

slide1 = pr1.slides.add_slide(slide1_register)

"main top placeholder"

title1 =  slide1.shapes.title
"placeholder = item in layout"
subtitle1 = slide1.placeholders[1]


"insert text"

title1.text = "ANALYSTRISING"
subtitle1.text = "Perf Manager"


#part 2
"create bullet point slide"
slide2_register = pr1.slide_layouts[1]
slide2 = pr1.slides.add_slide(slide2_register)

"Edit Bullet Point Slide"
title2 = slide2.shapes.title
title2.text = "Now for some bullet points"

bullet_point_box = slide2.shapes

bullet_points_lvl1 = bullet_point_box.placeholders[1]
bullet_points_lvl1.text = "donne ta perf du jour"

bullet_points_lvl2 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl2.text = "to"
bullet_points_lvl2.level = 1

bullet_points_lvl3 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl3.text = "my"
bullet_points_lvl3.level = 2

bullet_points_lvl4 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl4.text = "CHANNEL"
bullet_points_lvl4.level = 3


#slide 3
"create bullet point & picture slide"
slide3_register = pr1.slide_layouts[5]
slide3 = pr1.slides.add_slide(slide3_register)

"Edit Bullet Point & picture Slide"
title3 = slide3.shapes.title
title3.text = "Picture Time!"

"add image"

img1 = "lyon.jpg"

width = Inches(3)
height = Inches(3)


from_top = Inches(7.5) - Inches(3.75) - Inches(1.5)
from_left = Inches(13.33) - Inches(6.66) - Inches(1.5)




add_picture = slide3.shapes.add_picture(img1,from_left,from_top, width=Inches(3), height=Inches(3))

#Slide 4

"Register/Create Slide"
slide4_register = pr1.slide_layouts[5]
slide4 = pr1.slides.add_slide(slide4_register)

"slide 4 Title"

title4 = slide4.shapes.title
title4.text = "Shapework"

"Create Shapes"

#Shape 1

left1 = top1 = width1 = height1 = Inches(2)
add_shape1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left1,top1,width1,height1)

#Shape 2

left2 = Inches(6)
top2 = Inches(2)
width2 = height2 = Inches(2)

arrow1 = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,left2,top2,width2,height2)

#Change Arrow main colour

fill_arrow1 = arrow1.fill
fill_arrow1.solid()
fill_arrow1.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

#Rotate  Shape
arrow1.rotation =  90









pr1.save("AnalystRising_PPT_Tutorial.pptx")
filename = "AnalystRising_PPT_Tutorial.pptx"
opener = "open" if sys.platform == "darwin" else "xdg-open"
subprocess.call([opener, filename])